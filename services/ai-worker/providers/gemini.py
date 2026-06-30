"""
Google Gemini AI provider implementation.

Uses gemini-2.0-flash for cost-effective multimodal analysis.
Supports: PDF, images, DOCX, PPTX (text extracted before sending).
"""
from __future__ import annotations

import io
import json
import logging
from typing import Any

import google.generativeai as genai
from google.generativeai.types import GenerateContentResponse

from .base import AnalysisResult

logger = logging.getLogger(__name__)

_MODEL_NAME = "gemini-2.0-flash"

_ANALYSIS_PROMPT = """
Analyze the provided educational material and return a JSON object with these fields:

{
  "title": "Topografik chizmachilik",
  "description": "Plain-text summary of what this material covers (2-5 sentences)",
  "blurb": "One punchy sentence in Uzbek that makes a curious student want to read this",
  "tags": ["keyword1", "keyword2"],
  "authors": ["Familiya Ism Otashorasi"],
  "language": "O'zbek",
  "publish_year": 2023,
  "country": "O'zbekiston",
  "page_count": 148,
  "suggested_category_name": "Best matching academic category (e.g. Matematika, Tarix, Informatika)"
}

Rules:
- title: copy the EXACT title text as it visually appears in the document (largest/most prominent heading). Do NOT paraphrase, summarize, or generate a new title — copy it word for word. ONLY strip appended material-type descriptors that are NOT part of the title itself: "darslik", "o'quv qo'llanma", "o'quv uslubiy qo'llanma", "kitob", "majmua" and their combinations. For textbooks: the title is the subject/topic name — strip the material-type suffix. For articles/papers/journals: copy the full title as written, no matter how long. CORRECT (textbook): "Topografik chizmachilik". CORRECT (article): "Oliy ta'lim tizimida muhandislik grafikasi fanlarini qiyosiy o'qitilishining asosiy xususiyatlari". WRONG: "Topografik chizmachilik: Oliy o'quv yurtlari uchun darslik", "Chizmachilik (Topografik chizmachilik) o'quv qo'llanma".
- title, description, blurb, tags must be in Uzbek language.
- tags: exactly 4-6 lowercase Uzbek keywords relevant to the content.
- authors: extract ONLY the main author(s) from the document cover/header (not editors, reviewers, or series editors). Return full names in "Familiya Ism Otashorasi" format (e.g. "Raxmonjonov Xasan Aliyevich"). Empty array [] if not found.
- language: return EXACTLY ONE value from this list — O'zbek, Rus, Ingliz, Qoraqalpoq, Tojik, Qozoq, Arabcha, Nemis, Fransuz, Xitoy, Yapon, Koreys, Turk, Fors, Italyan, Ispan. If multilingual, pick the dominant language. null if not found.
- country: return EXACTLY ONE value from this list — O'zbekiston, Rossiya, AQSH, Buyuk Britaniya, Germaniya, Fransiya, Xitoy, Yaponiya, Janubiy Koreya, Hindiston, Turkiya, Qozog'iston, Qirg'iziston, Tojikiston, Turkmaniston, Ukraina, Belarus, Ozarbayjon, Gruziya, Kanada, Avstriya, Italiya. null if not found.
- blurb: one short, natural sentence in Uzbek — write like recommending the book to a classmate, not like an advertisement.
- publish_year: integer year (e.g. 2023) or null if not found.
- page_count: integer or null if not determinable.
- suggested_category_name: single category name, or null if unclear.
- Return ONLY the raw JSON object — no markdown fences, no explanation.
"""

# MIME types that Gemini inline_data does NOT support — extract text instead
_DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
_PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
_TEXT_EXTRACT_MIMES = _DOCX_MIMES | _PPTX_MIMES

_MAX_TEXT_CHARS = 8000


def _extract_docx_text(data: bytes) -> str:
    from docx import Document  # type: ignore[import-untyped]

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)[:_MAX_TEXT_CHARS]


def _extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation  # type: ignore[import-untyped]

    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)[:_MAX_TEXT_CHARS]


class GeminiProvider:
    """
    Gemini multimodal AI provider.
    Injects media bytes into a Gemini API request and parses structured JSON output.
    """

    def __init__(self, api_key: str) -> None:
        if not api_key:
            raise ValueError("GEMINI_API_KEY is required for GeminiProvider")
        genai.configure(api_key=api_key)
        self._model = genai.GenerativeModel(model_name=_MODEL_NAME)

    async def analyze(
        self,
        media_bytes: bytes,
        mime_type: str,
        prompt: str = _ANALYSIS_PROMPT,
    ) -> AnalysisResult:
        import asyncio

        effective_prompt = prompt or _ANALYSIS_PROMPT
        result = await asyncio.to_thread(self._run_sync, media_bytes, mime_type, effective_prompt)
        return result

    def _run_sync(
        self,
        media_bytes: bytes,
        mime_type: str,
        prompt: str,
    ) -> AnalysisResult:
        """Synchronous inner call — run inside asyncio.to_thread."""
        if mime_type in _TEXT_EXTRACT_MIMES:
            if mime_type in _DOCX_MIMES:
                extracted = _extract_docx_text(media_bytes)
                label = "Hujjat (DOCX)"
            else:
                extracted = _extract_pptx_text(media_bytes)
                label = "Taqdimot (PPTX)"

            if not extracted.strip():
                raise ValueError(f"{label} faylidan matn topilmadi yoki fayl bo'sh")

            logger.info("Extracted %d chars from %s", len(extracted), label)
            content_part: dict[str, Any] = {"text": f"[{label} matni]\n{extracted}"}
        else:
            content_part = {
                "inline_data": {
                    "mime_type": mime_type,
                    "data": media_bytes,
                }
            }

        response: GenerateContentResponse = self._model.generate_content(
            [content_part, prompt],
            generation_config={"temperature": 0.2, "max_output_tokens": 2048},
        )

        raw_text: str = response.text.strip()
        logger.debug("Gemini raw response: %s", raw_text[:800])

        return self._parse_response(raw_text)

    @staticmethod
    def _parse_response(raw_text: str) -> AnalysisResult:
        """Parse the JSON from the model response into AnalysisResult."""
        text = raw_text
        if text.startswith("```"):
            lines = text.split("\n")
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines and lines[-1].strip() == "```":
                lines = lines[:-1]
            text = "\n".join(lines)

        try:
            data: dict[str, Any] = json.loads(text)
        except json.JSONDecodeError as exc:
            logger.error(
                "Failed to parse Gemini JSON response: %s | raw: %s",
                exc,
                raw_text[:500],
            )
            raise ValueError(f"Invalid JSON from Gemini: {exc}") from exc

        publish_year = data.get("publish_year")
        if publish_year is not None:
            try:
                publish_year = int(publish_year)
                if not (1900 <= publish_year <= 2100):
                    publish_year = None
            except (TypeError, ValueError):
                publish_year = None

        page_count = data.get("page_count")
        if page_count is not None:
            try:
                page_count = int(page_count)
                if page_count < 1:
                    page_count = None
            except (TypeError, ValueError):
                page_count = None

        authors = data.get("authors") or []
        if not isinstance(authors, list):
            authors = []
        authors = [str(a) for a in authors if a]

        tags = data.get("tags") or []
        if not isinstance(tags, list):
            tags = []
        tags = [str(t) for t in tags if t]

        return AnalysisResult(
            title=data.get("title") or None,
            description=data.get("description") or None,
            blurb=data.get("blurb") or None,
            tags=tags,
            authors=authors,
            language=data.get("language") or None,
            publish_year=publish_year,
            country=data.get("country") or None,
            page_count=page_count,
            suggested_category_name=data.get("suggested_category_name") or None,
        )
