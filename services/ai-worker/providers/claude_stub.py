"""
Anthropic Claude provider — claude-sonnet-4-6.

PDF handling strategy:
  - Text-based PDFs: extract text from first 10 pages (PyMuPDF), send as plain text.
    This reduces token usage from ~200K to ~3K for a typical 200-page document.
  - Scanned PDFs (no extractable text): send as binary document block, max 10 MB.
  - Images: send as base64 image block.
  - DOCX/PPTX: extract text via python-docx / python-pptx, send as plain text.
"""
from __future__ import annotations

import base64
import io
import json
import logging
from typing import Any

import anthropic
import fitz  # PyMuPDF

from .base import AnalysisResult

logger = logging.getLogger(__name__)

_MODEL = "claude-sonnet-4-6"

_ANALYSIS_PROMPT = """
Analyze the provided educational material and return a JSON object with these fields:

{
  "title": "Concise, descriptive title (max 100 chars)",
  "description": "Plain-text summary of what this material covers (2-5 sentences)",
  "blurb": "Short marketing teaser that motivates a student to read this (1-2 sentences, engaging tone)",
  "tags": ["keyword1", "keyword2", "keyword3"],
  "authors": ["Author Name 1", "Author Name 2"],
  "language": "Language the document is written in (e.g. O'zbek, Rus, Ingliz)",
  "publish_year": 2023,
  "publish_place": "Toshkent",
  "country": "Country of publication (e.g. O'zbekiston, Rossiya)",
  "page_count": 148,
  "suggested_category_name": "Best matching academic category (e.g. Matematika, Tarix, Informatika)"
}

Rules:
- All text fields (title, description, blurb, tags) must be in Uzbek language.
- tags: 4-6 lowercase Uzbek keywords relevant to the content.
- authors: extract from the document cover/header; empty array [] if not found.
- language: detect from the document's content language; use Uzbek name of language.
- publish_year: integer year (e.g. 2023) or null if not found.
- publish_place: city/region where the material was published, as printed on the title page (e.g. "Toshkent — 2016" -> "Toshkent"), transliterated to Uzbek Latin script; null if not found.
- country: null if not found.
- page_count: integer or null if not determinable.
- suggested_category_name: single category name, or null if unclear.
- Return ONLY the raw JSON object — no markdown fences, no explanation.
"""

_SUPPORTED_IMAGE_TYPES = {"image/jpeg", "image/png", "image/gif", "image/webp"}

_DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
_PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}

_MAX_PAGES = 10
_MAX_CHARS = 8000
_MAX_BINARY_BYTES = 10 * 1024 * 1024  # 10 MB


def _extract_pdf_text(pdf_bytes: bytes) -> str:
    """Extract text from the first _MAX_PAGES pages of a PDF. Returns '' for scanned/image-only PDFs."""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    total_pages = len(doc)
    pages_to_read = min(_MAX_PAGES, total_pages)
    parts = [doc[i].get_text() for i in range(pages_to_read)]
    text = "\n\n".join(parts).strip()
    doc.close()
    extracted = text[:_MAX_CHARS]
    logger.info(
        "PDF text extraction: %d chars from %d/%d pages",
        len(extracted),
        pages_to_read,
        total_pages,
    )
    return extracted


def _extract_docx_text(data: bytes) -> str:
    from docx import Document  # type: ignore[import-untyped]

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)[:_MAX_CHARS]


def _extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation  # type: ignore[import-untyped]

    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)[:_MAX_CHARS]


def _parse_response(raw_text: str) -> AnalysisResult:
    text = raw_text
    if text.startswith("```"):
        lines = text.split("\n")
        if lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        text = "\n".join(lines)

    try:
        data: dict[str, Any] = json.loads(text)
    except json.JSONDecodeError as exc:
        logger.error("Failed to parse Claude JSON response: %s | raw: %s", exc, raw_text[:500])
        raise ValueError(f"Invalid JSON from Claude: {exc}") from exc

    publish_year = data.get("publish_year")
    if publish_year is not None:
        try:
            publish_year = int(publish_year)
            if not (1900 <= publish_year <= 2100):
                publish_year = None
        except (TypeError, ValueError):
            publish_year = None

    page_count = data.get("page_count")
    if page_count is not None:
        try:
            page_count = int(page_count)
            if page_count < 1:
                page_count = None
        except (TypeError, ValueError):
            page_count = None

    authors = data.get("authors") or []
    if not isinstance(authors, list):
        authors = []
    authors = [str(a) for a in authors if a]

    tags = data.get("tags") or []
    if not isinstance(tags, list):
        tags = []
    tags = [str(t) for t in tags if t]

    return AnalysisResult(
        title=data.get("title") or None,
        description=data.get("description") or None,
        blurb=data.get("blurb") or None,
        tags=tags,
        authors=authors,
        language=data.get("language") or None,
        publish_year=publish_year,
        publish_place=(str(data.get("publish_place")).strip() or None) if data.get("publish_place") else None,
        country=data.get("country") or None,
        page_count=page_count,
        suggested_category_name=data.get("suggested_category_name") or None,
    )


class ClaudeProvider:
    """Anthropic Claude Sonnet 4.6 provider for educational material analysis."""

    def __init__(self, api_key: str) -> None:
        if not api_key:
            raise ValueError("ANTHROPIC_API_KEY is required for ClaudeProvider")
        self._client = anthropic.AsyncAnthropic(api_key=api_key)

    async def analyze(
        self,
        media_bytes: bytes,
        mime_type: str,
        prompt: str = _ANALYSIS_PROMPT,
    ) -> AnalysisResult:
        effective_prompt = prompt or _ANALYSIS_PROMPT
        content: list[dict[str, Any]]

        if mime_type in _SUPPORTED_IMAGE_TYPES:
            data_b64 = base64.standard_b64encode(media_bytes).decode()
            content = [
                {"type": "image", "source": {"type": "base64", "media_type": mime_type, "data": data_b64}},
                {"type": "text", "text": effective_prompt},
            ]

        elif mime_type == "application/pdf":
            extracted_text = _extract_pdf_text(media_bytes)

            if extracted_text:
                combined = f"{effective_prompt}\n\n--- HUJJAT MATNI (birinchi {_MAX_PAGES} bet) ---\n{extracted_text}"
                content = [{"type": "text", "text": combined}]
            else:
                if len(media_bytes) > _MAX_BINARY_BYTES:
                    raise ValueError(
                        f"Skanerdan PDF juda katta: {len(media_bytes) / 1024 / 1024:.1f} MB "
                        f"(max {_MAX_BINARY_BYTES // 1024 // 1024} MB). Matn ajratib bo'lmadi."
                    )
                logger.warning(
                    "Scanned PDF detected (%d bytes), sending as binary document block",
                    len(media_bytes),
                )
                data_b64 = base64.standard_b64encode(media_bytes).decode()
                content = [
                    {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": data_b64}},
                    {"type": "text", "text": effective_prompt},
                ]

        elif mime_type in _DOCX_MIMES:
            extracted_text = _extract_docx_text(media_bytes)
            if not extracted_text.strip():
                raise ValueError("DOCX faylidan matn topilmadi yoki fayl bo'sh")
            logger.info("DOCX text extraction: %d chars", len(extracted_text))
            combined = f"{effective_prompt}\n\n--- HUJJAT MATNI (DOCX) ---\n{extracted_text}"
            content = [{"type": "text", "text": combined}]

        elif mime_type in _PPTX_MIMES:
            extracted_text = _extract_pptx_text(media_bytes)
            if not extracted_text.strip():
                raise ValueError("PPTX faylidan matn topilmadi yoki fayl bo'sh")
            logger.info("PPTX text extraction: %d chars", len(extracted_text))
            combined = f"{effective_prompt}\n\n--- TAQDIMOT MATNI (PPTX) ---\n{extracted_text}"
            content = [{"type": "text", "text": combined}]

        else:
            raise ValueError(f"Unsupported MIME type for Claude: {mime_type}")

        response = await self._client.messages.create(
            model=_MODEL,
            max_tokens=2048,
            messages=[{"role": "user", "content": content}],
        )

        raw_text = response.content[0].text.strip()
        logger.debug("Claude raw response: %s", raw_text[:800])
        logger.info(
            "Claude token usage: input=%d output=%d",
            response.usage.input_tokens,
            response.usage.output_tokens,
        )
        return _parse_response(raw_text)
